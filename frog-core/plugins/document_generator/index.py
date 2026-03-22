import os
import json

def execute(params, context):
    action = params.get("action")
    path = params.get("path")
    content = params.get("content")

    if not path:
        return {"status": "error", "message": "Missing 'path' parameter. Must be an absolute path."}
    if not content:
        return {"status": "error", "message": "Missing 'content' parameter. Must be a valid JSON string."}

    try:
        abs_path = os.path.abspath(path)
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)
        
        # Parse the JSON string
        try:
            if isinstance(content, str):
                data = json.loads(content)
            else:
                data = content
        except Exception:
            return {"status": "error", "message": "Failed to parse 'content' as JSON. Ensure it is properly escaped."}

        if action == "create_ppt":
            try:
                from pptx import Presentation
            except ImportError:
                return {
                    "status": "error", 
                    "message": "Missing Dependency: Please install python-pptx using `pip install python-pptx` (via shell_executor) before attempting to create a PPT."
                }
            
            prs = Presentation()
            
            for slide_data in data:
                # Use bullet slide layout
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]
                
                title_shape.text = slide_data.get("title", "Untitled Slide")
                tf = body_shape.text_frame
                
                content_items = slide_data.get("content", [])
                if isinstance(content_items, str):
                    content_items = [content_items]
                
                for i, item in enumerate(content_items):
                    if i == 0:
                        tf.text = str(item)
                    else:
                        p = tf.add_paragraph()
                        p.text = str(item)
                        
            prs.save(abs_path)
            return {"status": "success", "message": f"Successfully created Presentation (.pptx) at: {abs_path}"}
            
        elif action == "create_word":
            try:
                from docx import Document
            except ImportError:
                return {
                    "status": "error", 
                    "message": "Missing Dependency: Please install python-docx using `pip install python-docx` (via shell_executor) before attempting to create a Word doc."
                }
                
            doc = Document()
            doc.add_heading('Generated Document', 0)
            
            for section in data:
                heading = section.get("heading")
                if heading:
                    doc.add_heading(str(heading), level=1)
                
                paragraphs = section.get("paragraphs", [])
                if isinstance(paragraphs, str):
                    paragraphs = [paragraphs]
                    
                for para in paragraphs:
                    doc.add_paragraph(str(para))
                    
            doc.save(abs_path)
            return {"status": "success", "message": f"Successfully created Word Document (.docx) at: {abs_path}"}

        else:
            return {"status": "error", "message": "Unsupported action. Must be 'create_ppt' or 'create_word'."}
            
    except Exception as e:
        return {"status": "error", "message": f"Document generation failed: {str(e)}"}
